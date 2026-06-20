#!/usr/bin/env python3
"""Verify MyStand Parser Tools with generated local samples."""

from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = ROOT / "test-output"


def write_text(path: Path, text: str) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")
    return path


def make_dxf(path: Path) -> Path:
    return write_text(
        path,
        "\n".join(
            [
                "0",
                "SECTION",
                "2",
                "HEADER",
                "0",
                "ENDSEC",
                "0",
                "SECTION",
                "2",
                "ENTITIES",
                "0",
                "TEXT",
                "8",
                "MyStandLayer",
                "10",
                "0",
                "20",
                "0",
                "40",
                "2.5",
                "1",
                "My Stand DXF sample",
                "0",
                "ENDSEC",
                "0",
                "EOF",
                "",
            ]
        ),
    )


def try_make_xlsx(path: Path) -> Path | None:
    try:
        from openpyxl import Workbook
    except Exception:
        return None
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "业绩"
    sheet.append(["姓名", "收入", "成本"])
    sheet.append(["刚哥", 12000, 3000])
    workbook.save(path)
    return path


def try_make_docx(path: Path) -> Path | None:
    try:
        from docx import Document
    except Exception:
        return None
    doc = Document()
    doc.add_heading("My Stand 文档样例", level=1)
    doc.add_paragraph("这是给解析工具验证用的 DOCX。")
    doc.save(path)
    return path


def try_make_pptx(path: Path) -> Path | None:
    try:
        from pptx import Presentation
    except Exception:
        return None
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "My Stand PPT 样例"
    slide.placeholders[1].text = "解析工具验证"
    presentation.save(path)
    return path


def try_make_pdf(path: Path) -> Path | None:
    try:
        from reportlab.pdfgen import canvas
    except Exception:
        return None
    pdf = canvas.Canvas(str(path))
    pdf.drawString(72, 720, "My Stand PDF sample")
    pdf.drawString(72, 700, "Parser verification")
    pdf.save()
    return path


def try_make_png(path: Path) -> Path | None:
    if not shutil.which("tesseract"):
        return None
    try:
        from PIL import Image, ImageDraw
    except Exception:
        return None
    image = Image.new("RGB", (520, 120), color="white")
    draw = ImageDraw.Draw(image)
    draw.text((20, 45), "My Stand OCR sample", fill="black")
    image.save(path)
    return path


def make_samples(sample_dir: Path) -> tuple[list[Path], list[str]]:
    if sample_dir.exists():
        shutil.rmtree(sample_dir)
    sample_dir.mkdir(parents=True, exist_ok=True)
    samples = [
        write_text(sample_dir / "sample.md", "# My Stand\n\n这是 Markdown 样例。"),
        write_text(sample_dir / "sample.txt", "这是 TXT 样例。\n第二行内容。"),
        write_text(sample_dir / "sample.csv", "姓名,收入,成本\n刚哥,12000,3000\n"),
        write_text(sample_dir / "sample.json", json.dumps({"name": "My Stand", "ok": True}, ensure_ascii=False)),
        write_text(sample_dir / "sample.xml", "<root><title>My Stand XML</title></root>"),
        write_text(sample_dir / "sample.html", "<html><body><h1>My Stand HTML</h1><p>网页样例正文。</p></body></html>"),
        make_dxf(sample_dir / "sample.dxf"),
    ]
    with zipfile.ZipFile(sample_dir / "sample.zip", "w") as archive:
        archive.writestr("note.md", "# ZIP 内文件\n\n这是压缩包里的文本。")
        archive.writestr("table.csv", "字段,值\nA,1\n")
    samples.append(sample_dir / "sample.zip")

    skipped: list[str] = []
    for label, maker in [
        ("sample.xlsx", try_make_xlsx),
        ("sample.docx", try_make_docx),
        ("sample.pptx", try_make_pptx),
        ("sample.pdf", try_make_pdf),
        ("sample.png", try_make_png),
    ]:
        generated = maker(sample_dir / label)
        if generated:
            samples.append(generated)
        else:
            skipped.append(label)
    return samples, skipped


def run_parser(parser_command: str, input_uri: str, output_path: Path, timeout: int) -> dict:
    command = [parser_command, "--input", input_uri, "--output", str(output_path)]
    completed = subprocess.run(command, check=False, capture_output=True, text=True, timeout=timeout)
    payload = {}
    if output_path.exists():
        try:
            payload = json.loads(output_path.read_text(encoding="utf-8"))
        except Exception as exc:
            payload = {"errors": [f"invalid parser json: {exc}"]}
    return {
        "returncode": completed.returncode,
        "stdout": completed.stdout.strip(),
        "stderr": completed.stderr.strip(),
        "payload": payload,
    }


def verify_sample(parser_command: str, sample: Path, output_dir: Path, timeout: int) -> tuple[bool, str]:
    output_path = output_dir / f"{sample.stem}.json"
    result = run_parser(parser_command, str(sample), output_path, timeout)
    payload = result["payload"]
    markdown = str(payload.get("content", {}).get("markdown") or "")
    errors = payload.get("errors") or []
    tool = payload.get("tool") or ""
    ok = result["returncode"] == 0 and bool(markdown.strip()) and not errors
    detail = f"{sample.name}: tool={tool} markdown={len(markdown)}"
    if not ok:
        detail += f" returncode={result['returncode']} errors={errors} stderr={result['stderr'][:200]}"
    return ok, detail


def verify_wechat(parser_command: str, url: str, output_dir: Path, timeout: int) -> tuple[bool, str]:
    output_path = output_dir / "wechat_article.json"
    result = run_parser(parser_command, url, output_path, timeout)
    payload = result["payload"]
    markdown = str(payload.get("content", {}).get("markdown") or "")
    errors = payload.get("errors") or []
    title = payload.get("source", {}).get("title") or ""
    tool = payload.get("tool") or ""
    ok = result["returncode"] == 0 and bool(markdown.strip()) and not errors
    detail = f"wechat: tool={tool} title={title!r} markdown={len(markdown)}"
    if not ok:
        detail += f" returncode={result['returncode']} errors={errors} stderr={result['stderr'][:200]}"
    return ok, detail


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Verify MyStand Parser Tools.")
    parser.add_argument("--parser", default=str(ROOT / "bin" / "mystand-parser"))
    parser.add_argument("--output", default=str(DEFAULT_OUTPUT))
    parser.add_argument("--timeout", type=int, default=90)
    parser.add_argument("--wechat-url", default="")
    parser.add_argument("--require-wechat", action="store_true")
    args = parser.parse_args(argv)

    output_dir = Path(args.output).expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    samples, skipped = make_samples(output_dir / "generated-samples")

    failures: list[str] = []
    for sample in samples:
        ok, detail = verify_sample(args.parser, sample, output_dir, args.timeout)
        print(("ok " if ok else "FAIL ") + detail)
        if not ok:
            failures.append(detail)

    for item in skipped:
        print(f"skip {item}: generator dependency or system binary not available")

    if args.wechat_url:
        ok, detail = verify_wechat(args.parser, args.wechat_url, output_dir, args.timeout)
        print(("ok " if ok else "WARN ") + detail)
        if args.require_wechat and not ok:
            failures.append(detail)

    if failures:
        print("\nFAILURES")
        for failure in failures:
            print(f"- {failure}")
        return 1

    print(f"\nparser verification completed; output={output_dir}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
